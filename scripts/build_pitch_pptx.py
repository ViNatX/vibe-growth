"""
Convert the pitch deck PNGs into a .pptx file for Google Slides import.

Each slide PNG is placed full-bleed on a widescreen (16:9) slide.
"""

from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

FIG_DIR = Path(__file__).parent.parent / "figures" / "pitch"

prs = Presentation()
# Set 16:9 widescreen dimensions
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slide_files = [
    FIG_DIR / "pitch_slide1.png",
    FIG_DIR / "pitch_slide2.png",
    FIG_DIR / "pitch_slide3.png",
    FIG_DIR / "pitch_slide4.png",
    FIG_DIR / "pitch_slide5.png",
]

blank_layout = prs.slide_layouts[6]  # blank slide

for slide_path in slide_files:
    slide = prs.slides.add_slide(blank_layout)
    # Add image stretched to fill the entire slide
    slide.shapes.add_picture(
        str(slide_path),
        left=Emu(0),
        top=Emu(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )
    print(f"  Added slide: {slide_path.name}")

out_path = FIG_DIR / "ViNatX_Pitch.pptx"
prs.save(str(out_path))
print(f"\nPPTX saved: {out_path}")
